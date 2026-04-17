"""6-slide pptx (English, moderate text) for INFO7375 final."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "figures"
OUT = ROOT / "INFO7375_final_6tasks_en_v2.pptx"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(0x1F, 0x3B, 0x6E)
GREY = RGBColor(0x44, 0x44, 0x44)
ACCENT = RGBColor(0xC4, 0x4E, 0x52)


def add_title(slide, text, sub=None):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.1), Inches(1.2))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = text
    p.runs[0].font.size = Pt(32)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = NAVY
    if sub:
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.runs[0].font.size = Pt(15)
        p2.runs[0].font.color.rgb = GREY


def add_section(slide, left, top, width, label, body, lsize=15, bsize=16):
    tb = slide.shapes.add_textbox(left, top, width, Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.runs[0].font.size = Pt(lsize)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = ACCENT
    p.space_after = Pt(2)
    p2 = tf.add_paragraph()
    p2.text = body
    p2.runs[0].font.size = Pt(bsize)
    p2.runs[0].font.color.rgb = GREY


def add_bullets(slide, left, top, width, height, items, size=16):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (head, body) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i == 0:
            p.text = ""
        r = p.add_run()
        r.text = head + "  "
        r.font.bold = True
        r.font.size = Pt(size)
        r.font.color.rgb = NAVY
        r2 = p.add_run()
        r2.text = body
        r2.font.size = Pt(size)
        r2.font.color.rgb = GREY
        p.space_after = Pt(10)


def add_image(slide, path, left, top, width):
    slide.shapes.add_picture(str(path), left, top, width=width)


def add_footer(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.runs[0].font.size = Pt(10)
    p.runs[0].font.color.rgb = GREY
    p.runs[0].font.italic = True


blank = prs.slide_layouts[6]

# ---------- Slide 1: Part 1 ----------
s = prs.slides.add_slide(blank)
add_title(s, "Part 1 — TinyZero Countdown LoRA",
          "Qwen2.5-3B + PEFT LoRA + TRL GRPOTrainer")
add_bullets(s, Inches(0.8), Inches(1.9), Inches(11.8), Inches(5), [
    ("Task", "combine the given numbers to reach a target value"),
    ("Method", "freeze the 3B backbone, train a LoRA adapter (r=16) with GRPO"),
    ("Reward", "1 if the equation is correct and hits the target, else 0 — very sparse"),
    ("Pitfalls", "flash-attn wheels must match CUDA/Torch/Python; TRL GRPO API changed"),
    ("Outcome", "first ~100 steps stay near zero reward, then positive reward emerges"),
], size=18)

# ---------- Slide 2: Step 1 ----------
s = prs.slides.add_slide(blank)
add_title(s, "Step 1 — Qwen2.5-7B Baseline",
          "Full AgentFlow: Planner + Verifier + Executor + Generator")
add_image(s, FIG / "step1_qwen25_7b.png", Inches(0.4), Inches(1.7), width=Inches(7.0))
add_bullets(s, Inches(7.8), Inches(1.7), Inches(5.2), Inches(5), [
    ("Setup", "DashScope API, n=30/bench, max_steps=3, Serper search tool"),
    ("Highs", "hotpotqa 56.7% — multi-hop QA + search fits AgentFlow well"),
    ("Lows", "musique 10% — 4-hop reasoning is hard for a 7B model"),
    ("Insight", "the four-module design, not just model size, drives the score"),
    ("Pitfall", "DashScope JSON mode needs the word \"json\" in the prompt"),
], size=14)

# ---------- Slide 3: Step 2 ----------
s = prs.slides.add_slide(blank)
add_title(s, "Step 2 — Qwen3.5 Scaling (0.8B → 27B)",
          "Same full AgentFlow, five model sizes on five QA benchmarks")
add_image(s, FIG / "step2_scaling.png", Inches(0.4), Inches(1.7), width=Inches(7.6))
add_bullets(s, Inches(8.2), Inches(1.7), Inches(4.9), Inches(5), [
    ("Clean scaling", "bamboogle and 2wiki — each size step adds accuracy"),
    ("Saturation", "hotpotqa plateaus at 9B (~50%), 27B brings no gain"),
    ("Data artifact", "gaia 9B drop = 25/30 timeouts (killed at 300s)"),
    ("Hard task", "musique stays hard even at 27B (16.7%)"),
    ("Takeaway", "scaling helps but hits a ceiling on long-horizon tasks"),
], size=14)

# ---------- Slide 4: Step 3 ----------
s = prs.slides.add_slide(blank)
add_title(s, "Step 3 — BIRD Text-to-SQL",
          "Qwen3.5 + Serper + SQL executor, n = 22–30")
add_image(s, FIG / "step3_bird_scaling.png", Inches(0.4), Inches(1.7), width=Inches(7.0))
add_bullets(s, Inches(7.8), Inches(1.7), Inches(5.2), Inches(5), [
    ("Result", "non-monotonic: 4B peaks at 50%, 9B 36%, 27B 0%"),
    ("27B issue", "DashScope JSON errors — only 6 valid samples, stats unreliable"),
    ("9B issue", "extra explanation text around SQL → scorer cannot parse it"),
    ("Why no rerun", "API cost + sample sizes not comparable; document the gap"),
    ("Insight", "smaller models write cleaner SQL (less reasoning noise)"),
], size=14)

# ---------- Slide 5: Step 4 ----------
s = prs.slides.add_slide(blank)
add_title(s, "Step 4 — Flow-GRPO + LoRA Training",
          "Qwen3.5-0.8B, 30-step local smoke run (full runs on Modal)")
add_image(s, FIG / "step45_training_curves.png", Inches(0.4), Inches(1.7), width=Inches(7.8))
add_bullets(s, Inches(8.4), Inches(1.7), Inches(4.7), Inches(5), [
    ("New agent", "single model.generate() + regex — full AgentFlow too heavy in RL loop"),
    ("Reward", "acc_reward + think_format_reward, both 0/1"),
    ("Reward curve", "first fires at step 25 → LoRA starts learning"),
    ("KL", "stays small, updates well-controlled"),
    ("Entropy", "fluctuates but no collapse"),
], size=14)

# ---------- Slide 6: Step 5 ----------
s = prs.slides.add_slide(blank)
add_title(s, "Step 5 — Eval: Mind the Agent Mismatch",
          "Simplified agent (Step 4/5) vs full AgentFlow (Step 3)")
add_image(s, FIG / "step45_vs_step3.png", Inches(0.3), Inches(1.7), width=Inches(8.0))
add_bullets(s, Inches(8.6), Inches(1.7), Inches(4.5), Inches(5), [
    ("Red", "simplified agent — one generate + regex"),
    ("Blue", "full AgentFlow — 4 modules + tools"),
    ("Not comparable", "different pipelines, tools, and scoring"),
    ("BIRD 5% → 44%", "metric shift, not RL gain"),
    ("How to report", "trust training curves, not cross-agent deltas"),
], size=14)
add_footer(s, "Yubo Wang · INFO7375 Final · 2026-04-17")

prs.save(str(OUT))
print("wrote", OUT)
