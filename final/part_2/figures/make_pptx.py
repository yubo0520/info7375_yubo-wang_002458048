"""6-slide pptx for INFO7375 final: one slide per task."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "figures"
OUT = ROOT / "INFO7375_final_6tasks.pptx"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(0x1F, 0x3B, 0x6E)
GREY = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0xC4, 0x4E, 0x52)


def add_title(slide, text, sub=None):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.9))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = text
    p.runs[0].font.size = Pt(28)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = NAVY
    if sub:
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.runs[0].font.size = Pt(14)
        p2.runs[0].font.color.rgb = GREY


def add_bullets(slide, left, top, width, height, items, size=14):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (head, body) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r1 = p.add_run() if i == 0 else None
        if i == 0:
            p.text = ""
        r = p.add_run()
        r.text = head
        r.font.bold = True
        r.font.size = Pt(size)
        r.font.color.rgb = NAVY
        r2 = p.add_run()
        r2.text = "  " + body
        r2.font.size = Pt(size)
        r2.font.color.rgb = GREY
        p.space_after = Pt(8)


def add_image(slide, path, left, top, width=None, height=None):
    if width:
        slide.shapes.add_picture(str(path), left, top, width=width)
    else:
        slide.shapes.add_picture(str(path), left, top, height=height)


def add_footer(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.runs[0].font.size = Pt(10)
    p.runs[0].font.color.rgb = GREY
    p.runs[0].font.italic = True


blank = prs.slide_layouts[6]

# ---------- Slide 1: Part 1 TinyZero LoRA ----------
s = prs.slides.add_slide(blank)
add_title(s, "Part 1 — TinyZero Countdown LoRA",
          "Qwen2.5-3B-Instruct + PEFT LoRA + TRL GRPOTrainer")
add_bullets(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5), [
    ("目标", "复现 TinyZero：在 Countdown 任务（用给定数字凑目标数）上用 RL 微调一个小模型，\n     让它学会自己推理出正确的算式。"),
    ("方法", "base = Qwen2.5-3B-Instruct，冻结主干，加 LoRA（r=16），用 GRPO（group relative\n     policy optimization）在 ~500 条 prompts 上迭代，reward = 算式正确并等于目标。"),
    ("产出", "checkpoint 在 ./countdown_lora_3b/，训练脚本 train.py，Colab 版本 colab_run.ipynb。"),
    ("踩的坑", "flash-attn 轮子要对齐 CUDA/Torch/Python 版本；TRL 新版 GRPOTrainer 的 reward\n     签名变了，老教程的 prompt_template 对不上，照文档改。"),
    ("结论", "LoRA+GRPO 可以在 3B 量级跑起来，但 reward 比较稀疏（要完全正确才给分），\n     所以训练需要足够 rollout，否则前 100 步几乎都是 0 reward。"),
], size=15)

# ---------- Slide 2: Step 1 Qwen2.5-7B baseline ----------
s = prs.slides.add_slide(blank)
add_title(s, "Step 1 — Qwen2.5-7B 基线",
          "完整 AgentFlow（Planner + Verifier + Executor + Generator）")
add_image(s, FIG / "step1_qwen25_7b.png", Inches(0.4), Inches(1.5), width=Inches(6.5))
add_bullets(s, Inches(7.2), Inches(1.5), Inches(5.8), Inches(5.5), [
    ("设置", "DashScope API（OpenAI 兼容），n=30/bench，max_steps=3。\n     工具：Base_Generator + Serper_Search。"),
    ("结果", "hotpotqa 56.7% 最高（多跳 QA + 搜索最吃香），\n     musique 10% 最低（4 跳推理太难），\n     bamboogle/2wiki/gaia 处在中间。"),
    ("分析", "7B 开箱就能在简单 QA 上 50%+，说明 full AgentFlow 四模块\n     分工（planner 拆任务、verifier 过滤错误、executor 调工具）\n     的确有用，不只靠模型本身能力。"),
    ("踩的坑", "DashScope 的 response_format=json_object 要求 prompt 里含\n     \"json\" 字样，否则 400。后面 27b 跑 bird 就是踩到这个。"),
], size=13)

# ---------- Slide 3: Step 2 Qwen3.5 scaling ----------
s = prs.slides.add_slide(blank)
add_title(s, "Step 2 — Qwen3.5 0.8B → 27B scaling",
          "同一个 full AgentFlow，不同 size 的模型")
add_image(s, FIG / "step2_scaling.png", Inches(0.4), Inches(1.5), width=Inches(7.3))
add_bullets(s, Inches(8.0), Inches(1.5), Inches(5.0), Inches(5.5), [
    ("观察", "bamboogle 和 2wiki 呈现清晰的 scaling curve：\n     size 翻倍 → 准确率明显上台阶。"),
    ("意外", "hotpotqa 在 9B 就饱和了（~50%），27B 没再涨；\n     musique 很难，27B 也才 16.7%，不是 size 能解决的。"),
    ("gaia 的坑", "gaia 有 25/30 在 9B 跑 timeout（>300s 被 kill），\n     所以 9B 的 6.7% 是被 timeout 压低了，27B 回到 26.7%。"),
    ("结论", "scaling 有效但不是万能，4 跳推理和长时任务对\n     agent 框架提出更高要求。"),
], size=13)

# ---------- Slide 4: Step 3 BIRD ----------
s = prs.slides.add_slide(blank)
add_title(s, "Step 3 — BIRD Text-to-SQL",
          "Qwen3.5 + Serper 搜索 + SQL 执行器，n=22–30")
add_image(s, FIG / "step3_bird_scaling.png", Inches(0.4), Inches(1.5), width=Inches(6.8))
add_bullets(s, Inches(7.5), Inches(1.5), Inches(5.5), Inches(5.5), [
    ("结果反常", "scaling 不单调：4B 峰值 50%，9B 掉到 36%，\n     27B 只有 0%（n=6，3 个 timeout，3 个真实运行但答错）。"),
    ("原因 1（27B）", "DashScope 结构化输出踩坑，多次重试仍失败；\n     剩下 6 个样本太少，统计意义不强。"),
    ("原因 2（9B）", "output 文本包含 SQL 之外的解释文字，\n     calculate_score_bird.py 抽 SQL 失败。"),
    ("不补跑的原因", "DashScope 计费 + 时间成本，且报告里\n     已经有明确标注：n=6 vs n=22 不在同一量级。"),
    ("启示", "Text-to-SQL 对 prompt 格式非常敏感，小模型反而\n     输出更干净（less chain-of-thought noise）。"),
], size=12)

# ---------- Slide 5: Step 4 Flow-GRPO 训练 ----------
s = prs.slides.add_slide(blank)
add_title(s, "Step 4 — Flow-GRPO + LoRA 训练",
          "Qwen3.5-0.8B，本地 30 步冒烟，Modal 上跑全量")
add_image(s, FIG / "step45_training_curves.png", Inches(0.4), Inches(1.4), width=Inches(7.8))
add_bullets(s, Inches(8.4), Inches(1.4), Inches(4.6), Inches(5.7), [
    ("为啥另起一套 agent", "full AgentFlow 四模块 + API 调用太重，\n     放进 GRPO inner loop 跑不动（每 step 要几十次 API）。\n     task56_iso_new 改成单次 model.generate() + 正则解析标签。"),
    ("reward 设计", "acc_reward（答对）+ think_format_reward\n     （有 <think>/<answer> 标签），两者 0/1 加权。"),
    ("曲线怎么看", "reward 在 step 25 第一次非零 → LoRA 开始生效；\n     loss 对应尖峰；entropy 震荡但没塌陷；\n     kl 很小 → 更新幅度受控。"),
    ("坑", "冒烟跑只有 30 步 + 小 batch，reward 不稳是正常的；\n     全量 run 在 Modal 上（GPU 成本控制）。"),
], size=11)

# ---------- Slide 6: Step 5 eval & 为啥不和 Step 3 比 ----------
s = prs.slides.add_slide(blank)
add_title(s, "Step 5 — Eval & 架构不匹配的坑",
          "simplified agent（Step 4/5） vs full AgentFlow（Step 3）")
add_image(s, FIG / "step45_vs_step3.png", Inches(0.3), Inches(1.4), width=Inches(8.2))
add_bullets(s, Inches(8.8), Inches(1.4), Inches(4.3), Inches(5.7), [
    ("数字差异", "Step 5 在 2wiki 比 Step 3 高（0.30 vs 0.27），\n     但 bamboogle/hotpotqa/gaia 反而低。\n     bird 高很多（0.44 vs 0.05），主要因为\n     simplified agent 的 score_avg 更松。"),
    ("不能直接比较", "两套 agent 的推理流程不一样：\n     • Step 3：4 模块 + Serper + SQL 执行器\n     • Step 4/5：单次 generate + 正则抓答案\n     metric 的口径和工具访问都不同。"),
    ("正确的汇报方式", "训练稳定性看 reward/loss curve（已收敛），\n     不把 eval 数字减 Step 3 说成 \"RL 增益\"。"),
    ("教训", "RL 训练 loop 和推理 framework 解耦是工程常识，\n     但报告时要把口径讲清楚，否则 delta 没意义。"),
], size=11)

add_footer(s, "Yubo Wang · INFO7375 Final · 2026-04-17")

prs.save(str(OUT))
print("wrote", OUT)
